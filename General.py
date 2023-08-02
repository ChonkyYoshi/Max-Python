import docx
import pptx
from pathlib import Path
from zipfile import ZipFile
from lxml import etree
import openpyxl as xl
from shutil import make_archive, move, rmtree


def Unhide(File: Path,
           SkipRow: bool = False,
           SkipCol: bool = False,
           SkipSheet: bool = False,
           SkipShp: bool = False,
           SkipSld: bool = False,
           Overwrite: bool = False):

    match File.suffix:
        case '.docx' | '.docm':
            Doc = docx.Document(File)
            for index, par in enumerate(Doc.paragraphs):
                yield f'Paragraph {index + 1} of {len(Doc.Paragraphs)}'
                for run in par.runs:
                    run.font.hidden = False
            for index, table in enumerate(Doc.tables):
                yield f'Table {index + 1} of {len(Doc.tables)}'
                for row in table.rows:
                    for cell in row.cells:
                        for par in cell.paragraphs:
                            for run in par.runs:
                                run.font.hidden = False
            for section in Doc.sections:
                for par in section.header.paragraphs:
                    for run in par.runs:
                        run.font.hidden = False
                for par in section.footer.paragraphs:
                    for run in par.runs:
                        run.font.hidden = False
            if Overwrite:
                Doc.save(File.as_posix())
            else:
                Doc.save(f'{File.parent.as_posix()}/UNH_{File.name}')
        case '.xlsx' | '.xlsm':
            Sheets = list()
            wb = xl.load_workbook(File, read_only=True)
            for ws in wb.worksheets:
                if ws.sheet_state != 'visible':
                    Sheets.append(f'sheet{wb.get_index(ws) + 1}.xml')
            wb.close()
            Temp = Path(f'{File.parent.as_posix()}/Temp')
            Path.mkdir(Temp, exist_ok=True)
            ZipFile(File).extractall(Temp)
            for Sheetindex, Sheetfile in enumerate(
                    list(Temp.rglob('xl/worksheets/*.xml'))):
                Sheetcount = len(list(Temp.rglob('xl/worksheets/*.xml')))
                Sheetfile = Path(Sheetfile)
                Sheetxml = etree.parse(source=Sheetfile,
                                       parser=etree.XMLParser())
                if Sheetfile.name in Sheets and SkipSheet:
                    continue
                else:
                    wbfile = Path(f'{Temp.as_posix()}/xl/workbook.xml')
                    wbxml = etree.parse(source=wbfile,
                                        parser=etree.XMLParser())
                    for sheet in wbxml.xpath('//*[local-name()="sheet"]'):
                        sheet.set("state", "visible")
                if not SkipRow:
                    RowCount = len(list(
                        Sheetxml.xpath('//*[local-name()="row"]')))
                    for Rowindex, row in enumerate(
                            Sheetxml.xpath('//*[local-name()="row"]')):
                        yield f'sheet {Sheetindex + 1} of {Sheetcount}' +\
                            f'\nrow {Rowindex + 1} of {RowCount}'
                        row.set("hidden", "0")
                if not SkipCol:
                    ColCount = len(list(
                        Sheetxml.xpath('//*[local-name()="col"]')))
                    for ColIndex, col in enumerate(
                            Sheetxml.xpath('//*[local-name()="col"]')):
                        yield f'sheet {Sheetindex + 1} of {Sheetcount}' +\
                            f'\ncolumn {ColIndex + 1} of {ColCount}'
                        col.set('hidden', '0')
                with open(Sheetfile, 'wb') as f:
                    f.write(etree.tostring(Sheetxml))
            if Overwrite:
                new = Path(make_archive(File.stem, 'zip', root_dir=Temp))
                move(new, File.as_posix())
            else:
                new = Path(make_archive(f'UNH_{File.stem}',
                                        'zip', root_dir=Temp))
                move(new, f'{File.parent.as_posix()}/{new.stem}{File.suffix}')
            rmtree(Temp)
        case '.pptx' | '.pptm':
            Pres = pptx.Presentation(File)
            for index, slide in enumerate(Pres.slides):
                yield f'Slide {index} of {len(Pres.slides)}'
                if not SkipSld:
                    slide._element.set('show', '1')
                if not SkipShp:
                    for shape in slide.shapes:
                        hide(shape)
            if Overwrite:
                Pres.save(File)
            else:
                Pres.save(f'{File.parent.as_posix()}/UNH_{File.name}')


def hide(shape):

    match shape.shape_type:
        case 6:
            for sub in shape.shapes:
                hide(sub)
        case 13:
            shape._element.nvPicPr.cNvPr.set('hidden', '0')
        case 9:
            shape._element.nvCxnSpPr.cNvPr.set('hidden', '0')
        case 3 | 19 | 7:
            shape._element.nvGraphicFramePr.cNvPr.set('hidden', '0')
        case _:
            try:
                shape._element.nvSpPr.cNvPr.set('hidden', '0')
            except Exception:
                breakpoint()
