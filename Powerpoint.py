import pptx
from pathlib import Path
from regex import search, findall
import docx


def PPTSections(File: Path):

    PPT = pptx.Presentation(File)
    xml = str(PPT.part.blob)
    if search(r'(<p14:section name=")(.*?)(")', xml):
        doc = docx.Document()
        table = doc.add_table(rows=1, cols=2)
        r = table.cell(0, 0).paragraphs[0].add_run()
        r.text = 'Source'
        r.font.hidden = True
        r = table.cell(0, 1).paragraphs[0].add_run()
        r.text = 'Target'
        r.font.hidden = True
        count = len(findall(r'(<p14:section name=")(.*?)(")', xml))
        for index, sec in enumerate(findall(r'(<p14:section name=")(.*?)(")',
                                            xml)):
            yield f'section {index + 1} of {count}'
            row = table.add_row()
            r = row.cells[0].paragraphs[0].add_run()
            r.text = sec[1]
            r.font.hidden = True
            r = row.cells[1].paragraphs[0].add_run()
            r.text = sec[1]
        doc.save(f'{File.parent.as_posix()}/Section Titles_{File.name}.docx')


def NormalizeSpacing(File: Path, Overwrite: bool):

    PPT = pptx.Presentation(File)
    max = len(PPT.slides)
    for index, slide in enumerate(PPT.slides):
        yield f'Slide {index} of {max}'
        for shape in slide.shapes:
            Normalize(shape)
    if Overwrite:
        PPT.save(File.as_posix())
    else:
        PPT.save(f'{File.parent.as_posix()}/Fixed_{File.name}')


def Normalize(shape):

    if shape.shape_type == 6:
        for subshape in shape.shapes:
            Normalize(subshape)
    else:
        if shape.has_text_frame:
            TF = shape.text_frame
            for par in TF.paragraphs:
                for run in par.runs:
                    run.font._rPr.set('spc', '0')
